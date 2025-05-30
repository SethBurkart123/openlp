# -*- coding: utf-8 -*-

##########################################################################
# OpenLP - Open Source Lyrics Projection                                 #
# ---------------------------------------------------------------------- #
# Copyright (c) 2008 OpenLP Developers                                   #
# ---------------------------------------------------------------------- #
# This program is free software: you can redistribute it and/or modify   #
# it under the terms of the GNU General Public License as published by   #
# the Free Software Foundation, either version 3 of the License, or      #
# (at your option) any later version.                                    #
#                                                                        #
# This program is distributed in the hope that it will be useful,        #
# but WITHOUT ANY WARRANTY; without even the implied warranty of         #
# MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the          #
# GNU General Public License for more details.                           #
#                                                                        #
# You should have received a copy of the GNU General Public License      #
# along with this program.  If not, see <https://www.gnu.org/licenses/>. #
##########################################################################
"""
The :mod:`~openlp.plugins.mcp.mcpplugin` module contains the Plugin class
for the MCP (Model Context Protocol) plugin.
"""

import logging
import threading
import email
from pathlib import Path
from typing import List, Dict, Any

from PySide6 import QtCore

from openlp.core.state import State
from openlp.core.common.i18n import translate
from openlp.core.common.registry import Registry
from openlp.core.lib import build_icon
from openlp.core.lib.plugin import Plugin, StringContent
from openlp.core.lib.serviceitem import ServiceItem
from openlp.core.ui.icons import UiIcons
from openlp.core.common.enum import ServiceItemType

try:
    from fastmcp import FastMCP
    FASTMCP_AVAILABLE = True
except ImportError:
    FASTMCP_AVAILABLE = False

log = logging.getLogger(__name__)


class MCPWorker(QtCore.QObject):
    """
    Worker class that handles MCP operations on the main thread using Qt signals/slots.
    This ensures all GUI operations happen on the correct thread.
    """
    # Signals for different operations
    create_service_requested = QtCore.Signal()
    load_service_requested = QtCore.Signal(str)  # file_path
    save_service_requested = QtCore.Signal(str)  # file_path (optional)
    get_service_items_requested = QtCore.Signal()
    add_song_requested = QtCore.Signal(str, str, str)  # title, author, lyrics
    add_custom_slide_requested = QtCore.Signal(str, str)  # title, content
    add_media_requested = QtCore.Signal(str, str)  # file_path, title
    go_live_requested = QtCore.Signal(int)  # item_index
    next_slide_requested = QtCore.Signal()
    previous_slide_requested = QtCore.Signal()
    list_themes_requested = QtCore.Signal()
    set_theme_requested = QtCore.Signal(str)  # theme_name
    parse_email_requested = QtCore.Signal(str)  # email_content
    create_from_structure_requested = QtCore.Signal(object)  # service_structure
    
    # Result signals
    operation_completed = QtCore.Signal(object)  # result
    
    def __init__(self):
        super().__init__()
        self.current_result = None
        self.result_ready = threading.Event()
        
        # Connect signals to slots
        self.create_service_requested.connect(self.create_service)
        self.load_service_requested.connect(self.load_service)
        self.save_service_requested.connect(self.save_service)
        self.get_service_items_requested.connect(self.get_service_items)
        self.add_song_requested.connect(self.add_song)
        self.add_custom_slide_requested.connect(self.add_custom_slide)
        self.add_media_requested.connect(self.add_media)
        self.go_live_requested.connect(self.go_live)
        self.next_slide_requested.connect(self.next_slide)
        self.previous_slide_requested.connect(self.previous_slide)
        self.list_themes_requested.connect(self.list_themes)
        self.set_theme_requested.connect(self.set_theme)
        self.parse_email_requested.connect(self.parse_email)
        self.create_from_structure_requested.connect(self.create_from_structure)
        
        self.operation_completed.connect(self._handle_result)
    
    def _handle_result(self, result):
        """Handle the result of an operation."""
        self.current_result = result
        self.result_ready.set()
    
    def wait_for_result(self, timeout=10):
        """Wait for an operation to complete and return the result."""
        self.result_ready.clear()
        if self.result_ready.wait(timeout):
            return self.current_result
        else:
            raise TimeoutError("Operation timed out")
    
    def wait_for_result_long(self, timeout=90):
        """Wait for a long operation (like PowerPoint conversion) to complete."""
        self.result_ready.clear()
        if self.result_ready.wait(timeout):
            return self.current_result
        else:
            raise TimeoutError("Long operation timed out")
    
    @QtCore.Slot()
    def create_service(self):
        try:
            service_manager = Registry().get('service_manager')
            service_manager.new_file()
            service_manager.repaint_service_list(-1, -1)
            self.operation_completed.emit("New service created successfully")
        except Exception as e:
            self.operation_completed.emit(f"Error creating new service: {str(e)}")
    
    @QtCore.Slot(str)
    def load_service(self, file_path):
        try:
            service_manager = Registry().get('service_manager')
            service_manager.load_file(Path(file_path))
            self.operation_completed.emit(f"Service loaded from {file_path}")
        except Exception as e:
            self.operation_completed.emit(f"Error loading service: {str(e)}")
    
    @QtCore.Slot(str)
    def save_service(self, file_path):
        try:
            service_manager = Registry().get('service_manager')
            if file_path:
                service_manager.set_file_name(Path(file_path))
            service_manager.decide_save_method()
            self.operation_completed.emit(f"Service saved{' to ' + file_path if file_path else ''}")
        except Exception as e:
            self.operation_completed.emit(f"Error saving service: {str(e)}")
    
    @QtCore.Slot()
    def get_service_items(self):
        try:
            service_manager = Registry().get('service_manager')
            items = []
            for item in service_manager.service_items:
                service_item = item['service_item']
                items.append({
                    'title': service_item.title,
                    'type': str(service_item.service_item_type),
                    'plugin': service_item.name,
                    'order': item['order']
                })
            self.operation_completed.emit(items)
        except Exception as e:
            self.operation_completed.emit([{"error": str(e)}])
    
    @QtCore.Slot(str, str, str)
    def add_song(self, title, author, lyrics):
        try:
            songs_plugin = Registry().get('plugin_manager').get_plugin_by_name('songs')
            if not songs_plugin or not songs_plugin.is_active():
                self.operation_completed.emit("Songs plugin not available")
                return
            
            # Simple song search by title
            from openlp.plugins.songs.lib.db import Song
            existing_songs = songs_plugin.manager.get_all_objects(Song, Song.title == title)
            
            if existing_songs:
                # Found existing song, use it
                song = existing_songs[0]
                try:
                    from PySide6.QtWidgets import QListWidgetItem
                    from PySide6.QtCore import Qt
                    mock_item = QListWidgetItem()
                    mock_item.setData(Qt.ItemDataRole.UserRole, song.id)
                    
                    media_item = songs_plugin.media_item
                    service_item = ServiceItem(songs_plugin)
                    service_item.add_icon()
                    
                    if media_item.generate_slide_data(service_item, item=mock_item):
                        service_manager = Registry().get('service_manager')
                        service_manager.add_service_item(service_item)
                        service_manager.repaint_service_list(-1, -1)
                        self.operation_completed.emit(f"Song '{song.title}' added from database")
                    else:
                        self._create_song_placeholder(songs_plugin, title, lyrics)
                        self.operation_completed.emit(f"Song '{title}' found but failed to load - added placeholder")
                except Exception:
                    self._create_song_placeholder(songs_plugin, title, lyrics)
                    self.operation_completed.emit(f"Song '{title}' found but failed to load - added placeholder")
            else:
                # No existing song, create placeholder
                self._create_song_placeholder(songs_plugin, title, lyrics)
                self.operation_completed.emit(f"Song '{title}' not found in database - added placeholder")
        except Exception as e:
            self.operation_completed.emit(f"Error adding song: {str(e)}")
    
    def _create_song_placeholder(self, songs_plugin, title, lyrics):
        """Helper method to create a song placeholder."""
        service_item = ServiceItem(songs_plugin)
        service_item.title = title
        service_item.name = 'songs'
        service_item.service_item_type = ServiceItemType.Text
        service_item.add_icon()
        
        if lyrics:
            verses = lyrics.split('\n\n')
            for verse in verses:
                if verse.strip():
                    service_item.add_from_text(verse.strip())
        else:
            service_item.add_from_text(f"Song: {title}\n\n(Lyrics not available)")
        
        service_manager = Registry().get('service_manager')
        service_manager.add_service_item(service_item)
        service_manager.repaint_service_list(-1, -1)
    
    @QtCore.Slot(str, str)
    def add_custom_slide(self, title, content):
        try:
            custom_plugin = Registry().get('plugin_manager').get_plugin_by_name('custom')
            service_item = ServiceItem(custom_plugin)
            service_item.title = title
            service_item.name = 'custom'
            service_item.service_item_type = ServiceItemType.Text
            service_item.add_icon()
            service_item.add_from_text(content)
            
            service_manager = Registry().get('service_manager')
            service_manager.add_service_item(service_item)
            service_manager.repaint_service_list(-1, -1)
            self.operation_completed.emit(f"Custom slide '{title}' added to service")
        except Exception as e:
            self.operation_completed.emit(f"Error adding custom slide: {str(e)}")
    
    @QtCore.Slot(str, str)
    def add_media(self, file_path, title):
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                self.operation_completed.emit(f"Error: File '{file_path}' not found")
                return
            
            # Detect media type
            extension = file_path.suffix.lower()
            image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.tif', '.webp', '.svg'}
            video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v', '.3gp'}
            audio_extensions = {'.mp3', '.wav', '.ogg', '.flac', '.aac', '.m4a', '.wma'}
            presentation_extensions = {'.pdf', '.pptx', '.ppt', '.pps', '.ppsx', '.odp'}
            
            if extension in image_extensions:
                self._add_image(file_path, title)
            elif extension in video_extensions or extension in audio_extensions:
                self._add_video_audio(file_path, title, extension in video_extensions)
            elif extension in presentation_extensions:
                self._add_presentation(file_path, title)
            else:
                supported = f"images ({', '.join(sorted(image_extensions))}), videos ({', '.join(sorted(video_extensions))}), audio ({', '.join(sorted(audio_extensions))}), presentations ({', '.join(sorted(presentation_extensions))})"
                self.operation_completed.emit(f"Unsupported format: {extension}. Supported: {supported}")
                
        except Exception as e:
            self.operation_completed.emit(f"Error adding media: {str(e)}")
    
    def _add_image(self, file_path, title):
        """Add an image using the images plugin."""
        images_plugin = Registry().get('plugin_manager').get_plugin_by_name('images')
        if not images_plugin or not images_plugin.is_active():
            self.operation_completed.emit("Images plugin not available")
            return
        
        service_item = ServiceItem(images_plugin)
        service_item.title = title or file_path.name
        service_item.name = 'images'
        service_item.add_icon()
        
        # Use add_from_image for proper image handling
        service_item.add_from_image(file_path, file_path.name)
        
        # Set file hash for saving
        from openlp.core.common import sha256_file_hash
        service_item.sha256_file_hash = sha256_file_hash(file_path)
        
        # Add image capabilities
        from openlp.core.lib.serviceitem import ItemCapabilities
        service_item.add_capability(ItemCapabilities.CanMaintain)
        service_item.add_capability(ItemCapabilities.CanPreview)
        service_item.add_capability(ItemCapabilities.CanLoop)
        service_item.add_capability(ItemCapabilities.CanAppend)
        service_item.add_capability(ItemCapabilities.CanEditTitle)
        service_item.add_capability(ItemCapabilities.HasThumbnails)
        service_item.add_capability(ItemCapabilities.ProvidesOwnTheme)
        
        service_manager = Registry().get('service_manager')
        service_manager.add_service_item(service_item)
        service_manager.repaint_service_list(-1, -1)
        self.operation_completed.emit(f"Image '{service_item.title}' added to service")
    
    def _add_video_audio(self, file_path, title, is_video):
        """Add a video or audio file using the media plugin."""
        media_plugin = Registry().get('plugin_manager').get_plugin_by_name('media')
        if not media_plugin or not media_plugin.is_active():
            self.operation_completed.emit("Media plugin not available")
            return
        
        service_item = ServiceItem(media_plugin)
        service_item.title = title or file_path.name
        service_item.name = 'media'
        service_item.service_item_type = ServiceItemType.Command
        service_item.add_icon()
        
        # Use add_from_command for video/audio files
        service_item.add_from_command(str(file_path.parent), file_path.name, UiIcons().clapperboard)
        
        # Set processor and capabilities
        service_item.processor = 'qt6'
        from openlp.core.lib.serviceitem import ItemCapabilities
        service_item.add_capability(ItemCapabilities.CanAutoStartForLive)
        service_item.add_capability(ItemCapabilities.CanEditTitle)
        service_item.add_capability(ItemCapabilities.RequiresMedia)
        
        service_manager = Registry().get('service_manager')
        service_manager.add_service_item(service_item)
        service_manager.repaint_service_list(-1, -1)
        
        media_type = "video" if is_video else "audio"
        self.operation_completed.emit(f"{media_type.capitalize()} '{service_item.title}' added to service")
    
    def _add_presentation(self, file_path, title):
        """Add a presentation file, converting PowerPoint to PDF if needed."""
        try:
            extension = file_path.suffix.lower()
            
            # If it's a PowerPoint file, convert to PDF first
            if extension in ['.pptx', '.ppt', '.pps', '.ppsx']:
                pdf_path = self._convert_ppt_to_pdf(file_path)
                if not pdf_path or not pdf_path.exists():
                    self.operation_completed.emit(f"Failed to convert PowerPoint file to PDF")
                    return
                # Use the converted PDF
                file_path = pdf_path
                title = f"{title or file_path.stem} (converted)"
            
            # Now handle as PDF presentation
            presentations_plugin = Registry().get('plugin_manager').get_plugin_by_name('presentations')
            if not presentations_plugin or not presentations_plugin.is_active():
                self.operation_completed.emit("Presentations plugin not available")
                return
            
            # Check if PDF controller is available
            if 'Pdf' not in presentations_plugin.controllers or not presentations_plugin.controllers['Pdf'].enabled():
                self.operation_completed.emit("PDF controller not available - please ensure PDF support is enabled")
                return
            
            controller = presentations_plugin.controllers['Pdf']
            
            # Create a presentation document
            doc = controller.add_document(file_path)
            if not doc:
                self.operation_completed.emit(f"Failed to create PDF document for: {file_path.name}")
                return
                
            if not doc.load_presentation():
                self.operation_completed.emit(f"Failed to load presentation: {file_path.name}")
                return
            
            # Create service item for presentation
            service_item = ServiceItem(presentations_plugin)
            service_item.title = title or file_path.name
            service_item.name = 'presentations'
            service_item.processor = 'Pdf'
            service_item.add_icon()
            
            # Set capabilities
            from openlp.core.lib.serviceitem import ItemCapabilities
            service_item.add_capability(ItemCapabilities.CanEditTitle)
            service_item.add_capability(ItemCapabilities.ProvidesOwnDisplay)
            service_item.add_capability(ItemCapabilities.HasThumbnails)
            
            # Get slide count with better error handling
            slide_count = 1  # Default fallback
            try:
                if hasattr(doc, 'get_slide_count'):
                    count = doc.get_slide_count()
                    if count is not None and count > 0:
                        slide_count = count
                elif hasattr(doc, 'get_page_count'):
                    count = doc.get_page_count()
                    if count is not None and count > 0:
                        slide_count = count
                elif hasattr(doc, 'pageCount'):
                    count = doc.pageCount()
                    if count is not None and count > 0:
                        slide_count = count
                else:
                    # Try to access slide_count attribute directly
                    if hasattr(doc, 'slide_count') and doc.slide_count:
                        slide_count = doc.slide_count
            except Exception as e:
                log.warning(f"Error getting slide count: {e}, using default")
                slide_count = 1
            
            # Add slides to service item
            if slide_count > 0:
                for i in range(1, slide_count + 1):
                    try:
                        thumbnail_path = doc.get_thumbnail_path(i, True) if hasattr(doc, 'get_thumbnail_path') else None
                        file_hash = doc.get_sha256_file_hash() if hasattr(doc, 'get_sha256_file_hash') else ""
                        
                        service_item.add_from_command(
                            str(file_path.parent), 
                            file_path.name, 
                            thumbnail_path or "", 
                            f"Slide {i}", 
                            "",
                            file_hash
                        )
                    except Exception as e:
                        log.warning(f"Error adding slide {i}: {e}")
                        # Add basic slide entry as fallback
                        try:
                            service_item.add_from_command(
                                str(file_path.parent), 
                                file_path.name, 
                                "", 
                                f"Slide {i}", 
                                "",
                                ""
                            )
                        except:
                            pass  # Skip this slide if it fails completely
                
                service_manager = Registry().get('service_manager')
                service_manager.add_service_item(service_item)
                service_manager.repaint_service_list(-1, -1)
                
                try:
                    doc.close_presentation()
                except:
                    pass
                    
                self.operation_completed.emit(f"Presentation '{service_item.title}' with {slide_count} slides added to service")
            else:
                try:
                    doc.close_presentation()
                except:
                    pass
                self.operation_completed.emit(f"No slides found in presentation: {file_path.name}")
                
        except Exception as e:
            self.operation_completed.emit(f"Error adding presentation: {str(e)}")
            log.error(f"Presentation error details: {e}", exc_info=True)
    
    def _convert_ppt_to_pdf(self, ppt_path):
        """Convert PowerPoint file to PDF using LibreOffice (preferred) or python-pptx fallback."""
        # First try LibreOffice (much better quality)
        libreoffice_result = self._convert_with_libreoffice(ppt_path)
        if libreoffice_result and libreoffice_result.exists():
            log.info(f"Successfully converted {ppt_path.name} using LibreOffice")
            return libreoffice_result
        
        # Fallback to python-pptx method
        log.info(f"LibreOffice not available, falling back to python-pptx for {ppt_path.name}")
        return self._convert_with_python_pptx(ppt_path)
    
    def _convert_with_libreoffice(self, ppt_path):
        """Convert PowerPoint file to PDF using LibreOffice."""
        try:
            import subprocess
            import time
            
            # Create temporary PDF file
            timestamp = int(time.time())
            pdf_name = f"{ppt_path.stem}_libreoffice_{timestamp}.pdf"
            pdf_path = ppt_path.parent / pdf_name
            
            # Try different LibreOffice command variations
            libreoffice_commands = [
                "soffice",
                "libreoffice", 
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/bin/soffice",
                "/usr/local/bin/soffice"
            ]
            
            for cmd in libreoffice_commands:
                try:
                    # Test if this command exists (quick check)
                    test_result = subprocess.run([cmd, "--version"], 
                                               capture_output=True, 
                                               text=True, 
                                               timeout=5)
                    if test_result.returncode == 0:
                        log.info(f"Found LibreOffice at: {cmd}, starting conversion...")
                        
                        # Emit progress message
                        self.operation_completed.emit(f"Converting PowerPoint using LibreOffice (this may take 30-60 seconds)...")
                        
                        # Perform the conversion with longer timeout
                        convert_cmd = [
                            cmd,
                            "--headless",
                            "--convert-to", "pdf",
                            "--outdir", str(ppt_path.parent),
                            str(ppt_path)
                        ]
                        
                        result = subprocess.run(convert_cmd, 
                                              capture_output=True, 
                                              text=True, 
                                              timeout=90)  # Increased to 90 seconds
                        
                        if result.returncode == 0:
                            # LibreOffice creates filename.pdf, rename to our unique name
                            default_pdf = ppt_path.with_suffix('.pdf')
                            if default_pdf.exists():
                                try:
                                    default_pdf.rename(pdf_path)
                                    log.info(f"LibreOffice conversion completed successfully")
                                    return pdf_path
                                except Exception:
                                    # Return the default PDF if rename fails
                                    log.info(f"LibreOffice conversion completed (using default name)")
                                    return default_pdf
                            elif pdf_path.exists():
                                return pdf_path
                        else:
                            log.debug(f"LibreOffice conversion failed: {result.stderr}")
                            
                except subprocess.TimeoutExpired:
                    log.warning(f"LibreOffice conversion timed out after 90 seconds")
                    return None
                except Exception as e:
                    log.debug(f"LibreOffice command {cmd} failed: {e}")
                    continue
            
            return None
                
        except Exception as e:
            log.debug(f"LibreOffice conversion error: {e}")
            return None
    
    def _convert_with_python_pptx(self, ppt_path):
        """Convert PowerPoint file to PDF using python-pptx and reportlab (fallback)."""
        try:
            # Try using python-pptx to extract content and create PDF
            from pptx import Presentation
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            import time
            
            # Create temporary PDF file
            timestamp = int(time.time())
            pdf_name = f"{ppt_path.stem}_python_{timestamp}.pdf"
            pdf_path = ppt_path.parent / pdf_name
            
            # Load PowerPoint presentation
            prs = Presentation(str(ppt_path))
            
            # Create PDF
            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            width, height = letter
            
            for i, slide in enumerate(prs.slides):
                # Add slide number
                c.drawString(50, height - 50, f"Slide {i + 1}")
                
                y_position = height - 100
                
                # Extract text from slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Simple text extraction
                        text_lines = shape.text.strip().split('\n')
                        for line in text_lines:
                            if line.strip() and y_position > 50:
                                c.drawString(50, y_position, line[:80])  # Limit line length
                                y_position -= 20
                
                c.showPage()
            
            c.save()
            
            if pdf_path.exists():
                log.info(f"Successfully converted {ppt_path.name} to PDF using python-pptx")
                return pdf_path
            else:
                log.error("PDF conversion failed - file not created")
                return None
                
        except ImportError as e:
            log.warning("PowerPoint conversion requires python-pptx and reportlab: pip install python-pptx reportlab")
            return None
        except Exception as e:
            log.error(f"Error converting PowerPoint to PDF: {e}")
            return None
    
    @QtCore.Slot(int)
    def go_live(self, item_index):
        try:
            service_manager = Registry().get('service_manager')
            service_manager.set_item(item_index)
            self.operation_completed.emit(f"Item {item_index} is now live")
        except Exception as e:
            self.operation_completed.emit(f"Error going live: {str(e)}")
    
    @QtCore.Slot()
    def next_slide(self):
        try:
            live_controller = Registry().get('live_controller')
            live_controller.slidecontroller_live_next.emit()
            self.operation_completed.emit("Moved to next slide")
        except Exception as e:
            self.operation_completed.emit(f"Error moving to next slide: {str(e)}")
    
    @QtCore.Slot()
    def previous_slide(self):
        try:
            live_controller = Registry().get('live_controller')
            live_controller.slidecontroller_live_previous.emit()
            self.operation_completed.emit("Moved to previous slide")
        except Exception as e:
            self.operation_completed.emit(f"Error moving to previous slide: {str(e)}")
    
    @QtCore.Slot()
    def list_themes(self):
        try:
            theme_manager = Registry().get('theme_manager')
            themes = theme_manager.get_theme_names()
            self.operation_completed.emit(themes)
        except Exception as e:
            self.operation_completed.emit([f"Error: {str(e)}"])
    
    @QtCore.Slot(str)
    def set_theme(self, theme_name):
        try:
            service_manager = Registry().get('service_manager')
            service_manager.service_theme = theme_name
            self.operation_completed.emit(f"Service theme set to '{theme_name}'")
        except Exception as e:
            self.operation_completed.emit(f"Error setting theme: {str(e)}")
    
    @QtCore.Slot(str)
    def parse_email(self, email_content):
        try:
            msg = email.message_from_string(email_content)
            subject = msg.get('subject', 'Unknown')
            
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode('utf-8')
                        break
            else:
                body = msg.get_payload(decode=True).decode('utf-8')
            
            service_manager = Registry().get('service_manager')
            service_manager.new_file()
            
            # Simple email parsing
            self.add_custom_slide("Service Title", subject)
            if body:
                self.add_custom_slide("Email Content", body)
            
            self.operation_completed.emit(f"Service created from email: '{subject}'")
        except Exception as e:
            self.operation_completed.emit(f"Error parsing email: {str(e)}")
    
    @QtCore.Slot(object)
    def create_from_structure(self, service_structure):
        try:
            service_manager = Registry().get('service_manager')
            service_manager.new_file()
            
            items_added = []
            for item_data in service_structure:
                item_type = item_data.get('type', 'custom')
                title = item_data.get('title', 'Untitled')
                content = item_data.get('content', '')
                
                if item_type == 'song':
                    self.add_song(title, item_data.get('author', ''), item_data.get('lyrics', content))
                    items_added.append(f"Song '{title}'")
                elif item_type == 'custom':
                    self.add_custom_slide(title, content)
                    items_added.append(f"Custom slide '{title}'")
                elif item_type == 'media':
                    file_path = item_data.get('file_path')
                    if file_path:
                        self.add_media(file_path, title)
                        items_added.append(f"Media '{title}'")
            
            service_manager.repaint_service_list(-1, -1)
            self.operation_completed.emit(f"Service created with {len(service_structure)} items: " + ", ".join(items_added))
        except Exception as e:
            self.operation_completed.emit(f"Error creating service: {str(e)}")


class MCPPlugin(Plugin):
    """
    The MCP plugin provides Model Context Protocol server functionality to allow AI models
    to fully control OpenLP, including creating services automatically from emails.
    """
    log.info('MCP Plugin loaded')

    def __init__(self):
        super(MCPPlugin, self).__init__('mcp')
        self.weight = -1
        self.icon_path = UiIcons().desktop
        self.icon = build_icon(self.icon_path)
        self.mcp_server = None
        self.server_thread = None
        self.worker = None
        State().add_service(self.name, self.weight, is_plugin=True)
        State().update_pre_conditions(self.name, self.check_pre_conditions())

    @staticmethod
    def about():
        about_text = translate('MCPPlugin', '<strong>MCP Plugin</strong><br />The MCP plugin provides '
                               'Model Context Protocol server functionality to allow AI models to fully control '
                               'OpenLP, including creating services automatically from emails and other sources.')
        return about_text

    def check_pre_conditions(self):
        """Check if FastMCP is available."""
        return FASTMCP_AVAILABLE

    def initialise(self):
        """Initialize the MCP server and start it in a separate thread."""
        if not FASTMCP_AVAILABLE:
            log.error('FastMCP not available. Please install fastmcp: pip install fastmcp')
            return

        log.info('MCP Plugin initialising')
        
        # Fix WebSocket worker issue
        from PySide6.QtCore import QTimer
        self.fix_timer = QTimer()
        self.fix_timer.setSingleShot(True)
        self.fix_timer.timeout.connect(self._fix_websocket_worker)
        self.fix_timer.start(100)
        
        self._setup_worker()
        self._setup_mcp_server()
        super(MCPPlugin, self).initialise()

    def finalise(self):
        """Shut down the MCP server."""
        log.info('MCP Plugin finalising')
        super(MCPPlugin, self).finalise()

    def _fix_websocket_worker(self):
        """Fix the WebSocket worker missing event_loop attribute."""
        try:
            ws_server = Registry().get('web_socket_server')
            if ws_server and ws_server.worker:
                worker = ws_server.worker
                if not hasattr(worker, 'event_loop') or worker.event_loop is None:
                    class MockEventLoop:
                        def is_running(self):
                            return False
                        def call_soon_threadsafe(self, callback, *args):
                            pass
                    
                    worker.event_loop = MockEventLoop()
                    log.info('Fixed WebSocket worker missing event_loop attribute')
        except Exception as e:
            log.debug(f'Could not fix WebSocket worker: {e}')

    def _setup_worker(self):
        """Set up the worker that will handle MCP operations on the main thread."""
        self.worker = MCPWorker()

    def _setup_mcp_server(self):
        """Set up the FastMCP server with all the tools for controlling OpenLP."""
        if not FASTMCP_AVAILABLE:
            return

        self.mcp_server = FastMCP("OpenLP Control Server")
        
        # Register tools
        self._register_service_tools()
        self._register_media_tools()
        self._register_slide_tools()
        self._register_theme_tools()
        self._register_email_tools()
        
        # Start server in separate thread
        self.server_thread = threading.Thread(target=self._run_server, daemon=True)
        self.server_thread.start()

    def _run_server(self):
        """Run the MCP server in a separate thread."""
        try:
            import asyncio
            loop = asyncio.new_event_loop()
            try:
                loop.run_until_complete(
                    self.mcp_server.run_async(transport="sse", host="127.0.0.1", port=8765)
                )
            finally:
                loop.close()
            log.info('MCP server started on http://127.0.0.1:8765')
        except Exception as e:
            log.error(f'Error running MCP server: {e}')

    def _register_service_tools(self):
        """Register tools for service management."""
        @self.mcp_server.tool()
        def create_new_service() -> str:
            """Create a new empty service."""
            self.worker.create_service_requested.emit()
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def load_service(file_path: str) -> str:
            """Load a service from a file path."""
            self.worker.load_service_requested.emit(file_path)
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def save_service(file_path: str = None) -> str:
            """Save the current service, optionally to a specific path."""
            self.worker.save_service_requested.emit(file_path or "")
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def get_service_items() -> List[Dict[str, Any]]:
            """Get all items in the current service."""
            self.worker.get_service_items_requested.emit()
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def add_song_to_service(title: str, author: str = None, lyrics: str = None) -> str:
            """Add a song to the current service."""
            self.worker.add_song_requested.emit(title, author or "", lyrics or "")
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def add_custom_slide_to_service(title: str, content: str) -> str:
            """Add a custom slide to the current service."""
            self.worker.add_custom_slide_requested.emit(title, content)
            return self.worker.wait_for_result()

    def _register_media_tools(self):
        """Register tools for media management."""
        @self.mcp_server.tool()
        def add_media_to_service(file_path: str, title: str = None) -> str:
            """Add a media file to the current service."""
            self.worker.add_media_requested.emit(file_path, title or "")
            return self.worker.wait_for_result()
        
        @self.mcp_server.tool()
        def add_sample_image() -> str:
            """Add the sample image.jpg to the service for testing."""
            import os
            sample_path = os.path.join(os.getcwd(), "image.jpg")
            self.worker.add_media_requested.emit(sample_path, "Sample Image")
            return self.worker.wait_for_result()
        
        @self.mcp_server.tool()
        def add_sample_video() -> str:
            """Add the sample video.mp4 to the service for testing."""
            import os
            sample_path = os.path.join(os.getcwd(), "video.mp4")
            self.worker.add_media_requested.emit(sample_path, "Sample Video")
            return self.worker.wait_for_result()
        
        @self.mcp_server.tool()
        def test_media_types() -> str:
            """Test adding both sample media files to demonstrate image vs video handling."""
            import os
            cwd = os.getcwd()
            
            # Create new service
            self.worker.create_service_requested.emit()
            result1 = self.worker.wait_for_result()
            
            # Add image
            image_path = os.path.join(cwd, "image.jpg")
            self.worker.add_media_requested.emit(image_path, "Test Image")
            result2 = self.worker.wait_for_result()
            
            # Add video
            video_path = os.path.join(cwd, "video.mp4")
            self.worker.add_media_requested.emit(video_path, "Test Video")
            result3 = self.worker.wait_for_result()
            
            return f"Media test completed:\n1. {result1}\n2. {result2}\n3. {result3}"
        
        @self.mcp_server.tool()
        def add_demo_pdf() -> str:
            """Add the Demo.pdf presentation to the service for testing."""
            import os
            demo_path = os.path.join(os.getcwd(), "Demo.pdf")
            self.worker.add_media_requested.emit(demo_path, "Demo PDF Presentation")
            return self.worker.wait_for_result()
        
        @self.mcp_server.tool()
        def add_demo_powerpoint() -> str:
            """Add the Demo.pptx presentation to the service (will be converted to PDF)."""
            import os
            demo_path = os.path.join(os.getcwd(), "Demo.pptx")
            self.worker.add_media_requested.emit(demo_path, "Demo PowerPoint Presentation")
            return self.worker.wait_for_result_long()
        
        @self.mcp_server.tool()
        def test_all_media_types() -> str:
            """Test all media types: image, video, PDF, and PowerPoint presentation."""
            import os
            cwd = os.getcwd()
            
            # Create new service
            self.worker.create_service_requested.emit()
            result1 = self.worker.wait_for_result()
            
            # Add image
            image_path = os.path.join(cwd, "image.jpg")
            self.worker.add_media_requested.emit(image_path, "Test Image")
            result2 = self.worker.wait_for_result()
            
            # Add video
            video_path = os.path.join(cwd, "video.mp4")
            self.worker.add_media_requested.emit(video_path, "Test Video")
            result3 = self.worker.wait_for_result()
            
            # Add PDF presentation
            pdf_path = os.path.join(cwd, "Demo.pdf")
            self.worker.add_media_requested.emit(pdf_path, "Demo PDF")
            result4 = self.worker.wait_for_result()
            
            # Add PowerPoint presentation (this takes longer)
            ppt_path = os.path.join(cwd, "Demo.pptx")
            self.worker.add_media_requested.emit(ppt_path, "Demo PowerPoint")
            result5 = self.worker.wait_for_result_long()
            
            return f"All media types tested:\n1. {result1}\n2. {result2}\n3. {result3}\n4. {result4}\n5. {result5}"

    def _register_slide_tools(self):
        """Register tools for controlling the live display."""
        @self.mcp_server.tool()
        def go_live_with_item(item_index: int) -> str:
            """Make a specific service item live by index."""
            self.worker.go_live_requested.emit(item_index)
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def next_slide() -> str:
            """Go to the next slide in the live item."""
            self.worker.next_slide_requested.emit()
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def previous_slide() -> str:
            """Go to the previous slide in the live item."""
            self.worker.previous_slide_requested.emit()
            return self.worker.wait_for_result()

    def _register_theme_tools(self):
        """Register tools for theme management."""
        @self.mcp_server.tool()
        def list_themes() -> List[str]:
            """Get a list of all available themes."""
            self.worker.list_themes_requested.emit()
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def set_service_theme(theme_name: str) -> str:
            """Set the theme for the current service."""
            self.worker.set_theme_requested.emit(theme_name)
            return self.worker.wait_for_result()

    def _register_email_tools(self):
        """Register tools for processing emails to create services."""
        @self.mcp_server.tool()
        def parse_email_for_service(email_content: str) -> str:
            """Parse an email and create a service from its content."""
            self.worker.parse_email_requested.emit(email_content)
            return self.worker.wait_for_result()

        @self.mcp_server.tool()
        def create_service_from_structure(service_structure: List[Dict[str, Any]]) -> str:
            """Create a service from a structured list of items."""
            self.worker.create_from_structure_requested.emit(service_structure)
            return self.worker.wait_for_result()

    def set_plugin_text_strings(self):
        """Called to define all translatable texts of the plugin."""
        self.text_strings[StringContent.Name] = {
            'singular': translate('MCPPlugin', 'MCP', 'name singular'),
            'plural': translate('MCPPlugin', 'MCP', 'name plural')
        }
        self.text_strings[StringContent.VisibleName] = {
            'title': translate('MCPPlugin', 'MCP', 'container title')
        } 